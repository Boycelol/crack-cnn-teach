import streamlit as st
import numpy as np
import matplotlib.pyplot as plt
from PIL import Image
import sqlite3
from datetime import datetime
import hashlib
import cv2
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import time
import os
from io import BytesIO
import base64
from pptx import Presentation

# ==================== 页面配置 ====================
st.set_page_config(
    page_title="CNN裂缝识别教学系统",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ==================== 自定义CSS美化 ====================
st.markdown("""
<style>
    .main-title {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.5rem;
        border-radius: 15px;
        text-align: center;
        margin-bottom: 2rem;
    }
    .main-title h1 {
        color: white;
        margin: 0;
        font-size: 2.5rem;
    }
    .main-title p {
        color: rgba(255,255,255,0.9);
        margin: 0.5rem 0 0 0;
    }
    .card {
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
        padding: 1.5rem;
        border-radius: 15px;
        margin: 1rem 0;
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }
    .result-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.5rem;
        border-radius: 15px;
        color: white;
        text-align: center;
    }
    .metric-card {
        background: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
        box-shadow: 0 2px 10px rgba(0,0,0,0.05);
        transition: transform 0.3s;
    }
    .metric-card:hover {
        transform: translateY(-5px);
    }
    .quiz-card {
        background: #f8f9fa;
        padding: 1.5rem;
        border-radius: 15px;
        margin: 1rem 0;
        border-left: 5px solid #667eea;
    }
    .analysis-card {
        background: linear-gradient(135deg, #f0f4ff 0%, #e8edff 100%);
        padding: 1.5rem;
        border-radius: 15px;
        margin: 1rem 0;
    }
    .example-img {
        border-radius: 10px;
        box-shadow: 0 4px 10px rgba(0,0,0,0.1);
        transition: transform 0.3s;
    }
    .example-img:hover {
        transform: scale(1.02);
    }
    .step-card {
        background: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
        margin: 0.5rem;
    }
    .badge {
        display: inline-block;
        padding: 0.2rem 0.8rem;
        border-radius: 20px;
        font-size: 0.8rem;
        font-weight: bold;
    }
    .badge-cnn { background: #e74c3c; color: white; }
    .badge-pool { background: #2ecc71; color: white; }
    .badge-fc { background: #f39c12; color: white; }
</style>
""", unsafe_allow_html=True)

# ==================== 修复中文显示 ====================
plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'Arial Unicode MS']
plt.rcParams['axes.unicode_minus'] = False

# ==================== 初始化数据库（带迁移）====================
def init_db():
    conn = sqlite3.connect('demo_records.db')
    cursor = conn.cursor()
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS records (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT,
            image_name TEXT,
            crack_result TEXT,
            confidence REAL,
            crack_count INTEGER,
            timestamp TEXT,
            parameters TEXT
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS quiz_records (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT,
            question_id INTEGER,
            question TEXT,
            user_answer TEXT,
            correct_answer TEXT,
            is_correct INTEGER,
            timestamp TEXT
        )
    ''')
    
    cursor.execute("PRAGMA table_info(quiz_records)")
    existing_columns = [column[1] for column in cursor.fetchall()]
    
    if 'category' not in existing_columns:
        cursor.execute('ALTER TABLE quiz_records ADD COLUMN category TEXT')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS analysis_records (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT,
            analysis_text TEXT,
            score INTEGER,
            timestamp TEXT
        )
    ''')
    
    conn.commit()
    conn.close()

init_db()

# ==================== 题库（10道题）====================
QUESTION_BANK = [
    {
        "id": 1,
        "question": "CNN中，卷积层的主要作用是什么？",
        "options": ["A. 降低图像分辨率", "B. 提取图像的局部特征", "C. 压缩图像文件大小", "D. 增强图像亮度"],
        "correct": "B",
        "correct_text": "提取图像的局部特征",
        "explanation": "卷积层通过卷积核在图像上滑动，计算局部区域的加权和，从而提取边缘、纹理等局部特征。",
        "category": "卷积层",
        "learning_action": "观看「卷积原理演示」页面，理解卷积核滑动过程"
    },
    {
        "id": 2,
        "question": "池化层（Pooling Layer）的主要功能是什么？",
        "options": ["A. 增加图像细节", "B. 降低特征图维度，保留主要特征", "C. 改变图像颜色空间", "D. 扩大感受野"],
        "correct": "B",
        "correct_text": "降低特征图维度，保留主要特征",
        "explanation": "池化层通过最大池化或平均池化降低特征图尺寸，减少计算量，同时保留最重要的特征。",
        "category": "池化层",
        "learning_action": "学习最大池化和平均池化的区别"
    },
    {
        "id": 3,
        "question": "下列哪个是CNN中常用的激活函数？",
        "options": ["A. Sigmoid", "B. Tanh", "C. ReLU", "D. 以上都是"],
        "correct": "D",
        "correct_text": "以上都是",
        "explanation": "ReLU、Sigmoid、Tanh都是CNN中常用的激活函数，其中ReLU最常用因为它计算简单且能缓解梯度消失。",
        "category": "激活函数",
        "learning_action": "记住ReLU公式：f(x)=max(0,x)"
    },
    {
        "id": 4,
        "question": "在裂缝识别任务中，浅层卷积层通常提取什么特征？",
        "options": ["A. 裂缝的整体形状", "B. 混凝土颜色分布", "C. 裂缝的边缘和纹理", "D. 裂缝的分类结果"],
        "correct": "C",
        "correct_text": "裂缝的边缘和纹理",
        "explanation": "浅层卷积层提取边缘、纹理等低级特征，这些特征对裂缝的边缘检测非常重要。",
        "category": "裂缝识别",
        "learning_action": "上传裂缝图片，观察检测效果"
    },
    {
        "id": 5,
        "question": "以下哪个操作可以防止CNN过拟合？",
        "options": ["A. 增加网络深度", "B. Dropout", "C. 使用更大的卷积核", "D. 减少训练数据"],
        "correct": "B",
        "correct_text": "Dropout",
        "explanation": "Dropout随机丢弃部分神经元，是防止过拟合的有效方法。",
        "category": "训练技巧",
        "learning_action": "了解Dropout、Batch Normalization等正则化方法"
    },
    {
        "id": 6,
        "question": "CNN中，感受野（Receptive Field）指的是什么？",
        "options": ["A. 卷积核的大小", "B. 输出特征图的一个像素对应输入图像的区域大小", "C. 池化窗口的大小", "D. 全连接层的神经元数量"],
        "correct": "B",
        "correct_text": "输出特征图的一个像素对应输入图像的区域大小",
        "explanation": "感受野是指特征图上的一个像素点对应输入图像的区域大小，深层网络的感受野更大。",
        "category": "基本概念",
        "learning_action": "复习感受野、步长、填充等基础概念"
    },
    {
        "id": 7,
        "question": "在混凝土裂缝检测中，为什么CNN比传统方法更有效？",
        "options": ["A. CNN可以自动学习特征", "B. CNN计算速度更快", "C. CNN不需要训练数据", "D. CNN只能处理小图片"],
        "correct": "A",
        "correct_text": "CNN可以自动学习特征",
        "explanation": "CNN能够自动从数据中学习裂缝的边缘、纹理、形态等层次化特征，无需人工设计特征。",
        "category": "裂缝识别",
        "learning_action": "尝试用不同图片测试裂缝检测效果"
    },
    {
        "id": 8,
        "question": "下列哪个是卷积核的常见尺寸？",
        "options": ["A. 1×1", "B. 3×3", "C. 5×5", "D. 以上都是"],
        "correct": "D",
        "correct_text": "以上都是",
        "explanation": "1×1、3×3、5×5等都是常见的卷积核尺寸，3×3最为常用。",
        "category": "卷积层",
        "learning_action": "了解不同尺寸卷积核的作用"
    },
    {
        "id": 9,
        "question": "CNN中，步长（Stride）的作用是什么？",
        "options": ["A. 控制卷积核滑动的间隔", "B. 控制学习率", "C. 控制批量大小", "D. 控制迭代次数"],
        "correct": "A",
        "correct_text": "控制卷积核滑动的间隔",
        "explanation": "步长决定了卷积核每次移动的像素数，步长越大，输出特征图尺寸越小。",
        "category": "卷积层",
        "learning_action": "观看卷积动画，理解步长的影响"
    },
    {
        "id": 10,
        "question": "以下哪个网络是经典的CNN架构？",
        "options": ["A. RNN", "B. LSTM", "C. ResNet", "D. Transformer"],
        "correct": "C",
        "correct_text": "ResNet",
        "explanation": "ResNet（残差网络）是经典的CNN架构，通过残差连接解决了深层网络的梯度消失问题。",
        "category": "经典网络",
        "learning_action": "了解ResNet、VGG、Inception等经典架构"
    }
]

# ==================== 生成示例图片 ====================
def create_sample_crack_image():
    """生成带裂缝的示例图片"""
    img = np.ones((400, 600, 3), dtype=np.uint8) * 200
    cv2.line(img, (100, 50), (500, 350), (30, 30, 30), 6)
    cv2.line(img, (300, 200), (150, 380), (40, 40, 40), 4)
    cv2.line(img, (350, 220), (500, 380), (35, 35, 35), 3)
    noise = np.random.randint(180, 220, (400, 600, 3), dtype=np.uint8)
    img = cv2.addWeighted(img, 0.7, noise, 0.3, 0)
    return img

def create_sample_normal_image():
    """生成正常混凝土示例图片"""
    img = np.ones((400, 600, 3), dtype=np.uint8) * 190
    noise = np.random.randint(170, 210, (400, 600, 3), dtype=np.uint8)
    img = cv2.addWeighted(img, 0.5, noise, 0.5, 0)
    return img

# ==================== 智能分析 ====================
def generate_smart_analysis(quiz_results, student_id, total_score):
    correct_count = sum(1 for r in quiz_results if r['is_correct'])
    total = len(quiz_results)
    
    topic_stats = {}
    for r in quiz_results:
        topic = r['category']
        if topic not in topic_stats:
            topic_stats[topic] = {'correct': 0, 'total': 0, 'wrong_questions': []}
        topic_stats[topic]['total'] += 1
        if r['is_correct']:
            topic_stats[topic]['correct'] += 1
        else:
            topic_stats[topic]['wrong_questions'].append(r['question'][:50])
    
    weak_topics = [(t, v['correct']/v['total']*100, v['wrong_questions']) 
                   for t, v in topic_stats.items() if v['correct']/v['total']*100 < 60]
    weak_topics.sort(key=lambda x: x[1])
    
    if total_score >= 90:
        level = "🏆 卓越"
        level_desc = "你对CNN的理解非常深入！"
        icon = "🌟"
    elif total_score >= 75:
        level = "📈 良好"
        level_desc = "基础扎实，继续加油！"
        icon = "👍"
    elif total_score >= 60:
        level = "📚 合格"
        level_desc = "掌握了核心概念，需要加强细节理解"
        icon = "📖"
    else:
        level = "⚠️ 需加强"
        level_desc = "建议重新学习基础知识"
        icon = "💪"
    
    report = f"""
<div class="analysis-card">
    <h3>{icon} 总体评价</h3>
    <p style="font-size:1.2rem">你的测验得分为 <strong style="font-size:1.5rem">{total_score}分</strong>（{correct_count}/{total}）</p>
    <p>评级：<strong>{level}</strong> — {level_desc}</p>
</div>

### 📊 知识点掌握情况

| 知识点 | 正确率 | 状态 | 建议 |
|:---|:---|:---|:---|
"""
    for topic, stats in topic_stats.items():
        rate = stats['correct'] / stats['total'] * 100
        if rate >= 80:
            status = "✅ 优秀"
            suggestion = "可挑战进阶内容"
        elif rate >= 60:
            status = "📖 良好"
            suggestion = "适当复习巩固"
        else:
            status = "⚠️ 需加强"
            suggestion = "重点复习"
        report += f"| {topic} | {rate:.0f}% | {status} | {suggestion} |\n"
    
    report += "\n---\n\n### 🎯 薄弱知识点分析\n\n"
    
    if weak_topics:
        for topic, rate, wrong_qs in weak_topics:
            report += f"<div style='background:#fff3f3; padding:10px; border-radius:10px; margin:10px 0;'>"
            report += f"<strong>🔴 {topic}</strong>（正确率 {rate:.0f}%）<br>"
            if wrong_qs:
                report += f"📝 错题示例：{wrong_qs[0]}<br>"
            report += f"</div><br>"
    else:
        report += "🎉 恭喜！所有知识点掌握良好！\n"
    
    report += """
---
### 💡 个性化学习建议

<div class="card">
    <ul>
"""
    
    suggestions = []
    if total_score < 60:
        suggestions.append("📖 <strong>系统复习</strong>：重新学习CNN结构图解页面")
        suggestions.append("🎬 <strong>观看演示</strong>：多看卷积动画演示")
    elif total_score < 80:
        suggestions.append("🔬 <strong>实践练习</strong>：上传不同裂缝图片进行检测")
        suggestions.append("📝 <strong>错题重做</strong>：仔细阅读错题解析")
    else:
        suggestions.append("🚀 <strong>进阶挑战</strong>：学习ResNet、DenseNet等架构")
    
    suggestions.append("🔍 <strong>裂缝检测</strong>：调节灵敏度参数观察效果")
    suggestions.append("📊 <strong>持续进步</strong>：定期重新测验，跟踪学习进步")
    
    for s in suggestions[:5]:
        report += f"        <li>{s}</li>\n"
    
    report += """
    </ul>
</div>

---
### 📚 推荐学习资源

<div class="card">
    <ul>
        <li>📺 李沐《动手学深度学习》- CNN章节</li>
        <li>🎓 吴恩达《Deep Learning Specialization》</li>
        <li>🏆 Kaggle 混凝土裂缝检测比赛</li>
        <li>📄 U-Net、ResNet 经典论文</li>
    </ul>
</div>
"""
    
    conn = sqlite3.connect('demo_records.db')
    conn.execute(
        'INSERT INTO analysis_records (session_id, analysis_text, score, timestamp) VALUES (?,?,?,?)',
        (student_id, report, total_score, datetime.now().isoformat())
    )
    conn.commit()
    conn.close()
    
    return report

# ==================== 会话ID ====================
def get_session_id():
    if 'session_id' not in st.session_state:
        st.session_state.session_id = hashlib.md5(
            str(datetime.now()).encode()
        ).hexdigest()[:8]
    return st.session_state.session_id

# ==================== 学习进度管理 ====================
def update_learning_progress(page_visited):
    if 'learning_progress' not in st.session_state:
        st.session_state.learning_progress = {
            'cnn_structure': False,
            'conv_demo': False,
            'crack_detection': False
        }
    if page_visited in st.session_state.learning_progress:
        st.session_state.learning_progress[page_visited] = True

def get_learning_progress():
    if 'learning_progress' not in st.session_state:
        return 0
    completed = sum(st.session_state.learning_progress.values())
    return int(completed / 3 * 100)

# ==================== CNN可视化函数 ====================
def draw_cnn_diagram():
    fig, ax = plt.subplots(figsize=(12, 4), facecolor='none')
    ax.set_facecolor('none')
    
    layers = [
        {"name": "输入图像", "x": 0, "color": "#3498db", "size": 0.45},
        {"name": "卷积层1", "x": 2, "color": "#e74c3c", "size": 0.45},
        {"name": "池化层1", "x": 4, "color": "#2ecc71", "size": 0.45},
        {"name": "卷积层2", "x": 6, "color": "#e74c3c", "size": 0.45},
        {"name": "池化层2", "x": 8, "color": "#2ecc71", "size": 0.45},
        {"name": "全连接层", "x": 10, "color": "#f39c12", "size": 0.45},
        {"name": "输出层", "x": 12, "color": "#9b59b6", "size": 0.45}
    ]
    
    for i in range(len(layers)-1):
        ax.annotate('', xy=(layers[i+1]['x'], 0), xytext=(layers[i]['x'], 0),
                   arrowprops=dict(arrowstyle='->', color='#888', lw=2))
    
    for layer in layers:
        circle = plt.Circle((layer['x'], 0), layer['size'], color=layer['color'], alpha=0.85)
        ax.add_patch(circle)
        ax.text(layer['x'], 0, layer['name'], ha='center', va='center', 
               fontsize=9, color='white', weight='bold')
    
    ax.set_xlim(-1, 13)
    ax.set_ylim(-1, 1)
    ax.axis('off')
    ax.set_title("卷积神经网络（CNN）结构图", fontsize=14, weight='bold', pad=20)
    
    plt.tight_layout()
    return fig

def draw_conv_demo():
    from scipy.signal import convolve2d
    fig, axes = plt.subplots(1, 3, figsize=(12, 3.5))
    
    img = np.zeros((50, 50))
    img[20:30, 15:35] = 1
    img = img + np.random.rand(50, 50) * 0.15
    axes[0].imshow(img, cmap='gray')
    axes[0].set_title('① 输入图像', fontsize=11)
    axes[0].axis('off')
    
    kernel = np.array([[-1, -1, -1], [-1, 8, -1], [-1, -1, -1]])
    im = axes[1].imshow(kernel, cmap='RdBu', vmin=-1, vmax=1)
    axes[1].set_title('② 3×3 卷积核', fontsize=11)
    axes[1].axis('off')
    
    output = convolve2d(img, kernel, mode='same')
    axes[2].imshow(output, cmap='hot')
    axes[2].set_title('③ 输出特征图', fontsize=11)
    axes[2].axis('off')
    
    plt.suptitle("卷积操作原理演示", fontsize=14, weight='bold')
    plt.tight_layout()
    return fig

def draw_relu_graph():
    """绘制ReLU函数图"""
    fig, ax = plt.subplots(figsize=(8, 4))
    x = np.linspace(-3, 3, 100)
    y = np.maximum(0, x)
    ax.plot(x, y, 'b-', linewidth=2, color='#e74c3c')
    ax.fill_between(x, 0, y, alpha=0.3, color='#e74c3c')
    ax.axhline(y=0, color='k', linestyle='-', alpha=0.3)
    ax.axvline(x=0, color='k', linestyle='-', alpha=0.3)
    ax.set_title('ReLU激活函数: f(x) = max(0, x)', fontsize=12)
    ax.set_xlabel('x')
    ax.set_ylabel('f(x)')
    ax.grid(True, alpha=0.3)
    plt.tight_layout()
    return fig

def draw_pooling_demo():
    """绘制池化演示图"""
    fig, axes = plt.subplots(1, 2, figsize=(10, 4))
    
    # 输入特征图
    input_data = np.array([
        [1, 3, 2, 4],
        [5, 1, 6, 8],
        [2, 0, 4, 3],
        [7, 2, 5, 9]
    ])
    axes[0].imshow(input_data, cmap='Blues', vmin=0, vmax=9)
    for i in range(4):
        for j in range(4):
            axes[0].text(j, i, str(input_data[i, j]), ha='center', va='center', fontsize=12)
    axes[0].set_title('输入特征图 (4×4)', fontsize=11)
    axes[0].axis('off')
    
    # 最大池化输出
    output = np.array([
        [5, 8],
        [7, 9]
    ])
    axes[1].imshow(output, cmap='Reds', vmin=0, vmax=9)
    for i in range(2):
        for j in range(2):
            axes[1].text(j, i, str(output[i, j]), ha='center', va='center', fontsize=16, weight='bold')
    axes[1].set_title('最大池化输出 (2×2)', fontsize=11)
    axes[1].axis('off')
    
    plt.suptitle("最大池化 (Max Pooling) 演示", fontsize=14, weight='bold')
    plt.tight_layout()
    return fig

# ==================== 裂缝检测函数 ====================
def detect_cracks(image, sensitivity=5, min_length=50, connect_distance=20):
    if isinstance(image, Image.Image):
        image = np.array(image)
    
    if len(image.shape) == 3:
        gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
    else:
        gray = image.copy()
    
    if len(image.shape) == 2:
        result_img = cv2.cvtColor(image, cv2.COLOR_GRAY2RGB)
    elif len(image.shape) == 3 and image.shape[2] == 4:
        result_img = cv2.cvtColor(image, cv2.COLOR_RGBA2RGB)
    else:
        result_img = image.copy()
    
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    enhanced = clahe.apply(gray)
    
    threshold_value = max(10, 50 - (sensitivity - 1) * 4)
    _, binary = cv2.threshold(enhanced, threshold_value, 255, cv2.THRESH_BINARY_INV)
    
    kernel = np.ones((connect_distance // 5 + 1, connect_distance // 5 + 1), np.uint8)
    binary = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel, iterations=2)
    
    kernel_erase = np.ones((2, 2), np.uint8)
    binary = cv2.morphologyEx(binary, cv2.MORPH_OPEN, kernel_erase, iterations=1)
    
    contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    crack_count = 0
    crack_info = []
    
    for contour in contours:
        length = cv2.arcLength(contour, False)
        area = cv2.contourArea(contour)
        x, y, bw, bh = cv2.boundingRect(contour)
        
        if bw > bh:
            aspect_ratio = bw / (bh + 0.01)
        else:
            aspect_ratio = bh / (bw + 0.01)
        
        is_crack = False
        if length > min_length:
            if aspect_ratio > 2.0:
                is_crack = True
            elif area > 100 and aspect_ratio > 1.5:
                is_crack = True
            elif sensitivity > 7 and length > min_length * 0.7:
                is_crack = True
        
        if is_crack:
            crack_count += 1
            crack_info.append({'length': length, 'area': area})
            cv2.rectangle(result_img, (x, y), (x + bw, y + bh), (255, 0, 0), 3)
    
    confidence = min(0.95, 0.3 + crack_count * 0.2) if crack_count > 0 else 0.0
    
    return {
        'has_crack': crack_count > 0,
        'crack_count': crack_count,
        'confidence': confidence,
        'result_image': result_img,
        'crack_info': crack_info
    }

# ==================== 统计分析 ====================
def get_quiz_statistics(session_id):
    conn = sqlite3.connect('demo_records.db')
    df = pd.read_sql_query(
        'SELECT * FROM quiz_records WHERE session_id = ?',
        conn, params=(session_id,)
    )
    conn.close()
    
    if df.empty:
        return None
    
    total = len(df)
    correct = df['is_correct'].sum()
    score = int(correct / total * 100) if total > 0 else 0
    
    topic_stats = {}
    for _, row in df.iterrows():
        topic = row['category'] if 'category' in row else '其他'
        if topic not in topic_stats:
            topic_stats[topic] = {'correct': 0, 'total': 0}
        topic_stats[topic]['total'] += 1
        if row['is_correct']:
            topic_stats[topic]['correct'] += 1
    
    return {
        'total': total,
        'correct': int(correct),
        'score': score,
        'topic_stats': topic_stats,
        'records': df.to_dict('records')
    }

# ==================== 主界面 ====================
def main():
    st.markdown("""
    <div class="main-title">
        <h1>🔬 混凝土裂缝智能识别教学系统</h1>
        <p>基于CNN的裂缝检测 | 互动式教学 | 智能学情分析</p>
    </div>
    """, unsafe_allow_html=True)
    
    with st.sidebar:
        st.markdown("### 📚 导航菜单")
        
        selected = st.radio(
            "选择页面",
            ["📖 CNN结构学习", "🎬 卷积原理演示", "🔍 裂缝检测实战", "📝 智能测验", "📊 学情分析"],
            label_visibility="collapsed",
            key="nav_menu"
        )
        
        st.markdown("---")
        
        st.markdown("### 📈 学习进度")
        progress = get_learning_progress()
        st.progress(progress / 100)
        st.caption(f"已完成 {progress}%")
        
        st.markdown("---")
        
        if st.button("📊 预览教学PPT", use_container_width=True):
            try:
                prs = Presentation("C6.pptx")
                st.success(f"PPT共 {len(prs.slides)} 页，可下载到本地完整查看")
                # 这里可以循环读取每页文字展示简易预览
                for idx, slide in enumerate(prs.slides, 1):
                    text = ""
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text + "\n"
                    st.subheader(f"第{idx}页")
                    st.text(text if text else "当前页面无文字内容")
            except Exception as e:
                st.error(f"PPT读取失败：{str(e)}")
        
        st.markdown("---")
        st.caption(f"🆔 会话ID: {get_session_id()}")
        st.caption(f"📅 {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    
    page_map = {
        "📖 CNN结构学习": "cnn_structure",
        "🎬 卷积原理演示": "conv_demo",
        "🔍 裂缝检测实战": "crack_detection"
    }
    if selected in page_map:
        update_learning_progress(page_map[selected])
    
    # ==================== 页面1：CNN结构学习 ====================
    if selected == "📖 CNN结构学习":
        st.header("📖 卷积神经网络结构学习")
        
        # 介绍卡片
        st.markdown("""
        <div class="card">
            <p>卷积神经网络（CNN）是深度学习的核心算法之一，特别适合处理图像数据。
            在混凝土裂缝检测中，CNN能够自动学习从边缘纹理到裂缝形态的层次化特征。</p>
        </div>
        """, unsafe_allow_html=True)
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.markdown("### 🏗️ CNN网络结构")
            st.pyplot(draw_cnn_diagram())
        
        with col2:
            st.markdown("### 🎯 CNN核心组件")
            
            with st.expander("📌 卷积层", expanded=True):
                st.markdown("""
                <span class="badge badge-cnn">核心</span>
                - **作用**：提取图像局部特征
                - **原理**：卷积核滑动计算
                - **裂缝应用**：检测裂缝边缘纹理
                """, unsafe_allow_html=True)
            
            with st.expander("📌 池化层"):
                st.markdown("""
                <span class="badge badge-pool">压缩</span>
                - **作用**：降维压缩特征
                - **原理**：最大/平均池化
                - **裂缝应用**：保留主要裂缝特征
                """, unsafe_allow_html=True)
            
            with st.expander("📌 全连接层"):
                st.markdown("""
                <span class="badge badge-fc">分类</span>
                - **作用**：整合所有特征
                - **原理**：加权求和+激活
                - **裂缝应用**：裂缝/正常分类
                """, unsafe_allow_html=True)
        
        # CNN四步流程
        st.markdown("---")
        st.markdown("### 🔬 CNN如何识别混凝土裂缝？")
        
        col_a, col_b, col_c, col_d = st.columns(4)
        with col_a:
            st.markdown("""
            <div class="metric-card">
                <h3 style="color:#3498db">📥 输入</h3>
                <p>混凝土图像</p>
                <small>224×224×3</small>
            </div>
            """, unsafe_allow_html=True)
        with col_b:
            st.markdown("""
            <div class="metric-card">
                <h3 style="color:#e74c3c">🔍 卷积</h3>
                <p>提取边缘纹理</p>
                <small>64个特征图</small>
            </div>
            """, unsafe_allow_html=True)
        with col_c:
            st.markdown("""
            <div class="metric-card">
                <h3 style="color:#2ecc71">📊 池化</h3>
                <p>降维压缩特征</p>
                <small>保留主要特征</small>
            </div>
            """, unsafe_allow_html=True)
        with col_d:
            st.markdown("""
            <div class="metric-card">
                <h3 style="color:#9b59b6">✅ 输出</h3>
                <p>裂缝/正常判断</p>
                <small>二分类结果</small>
            </div>
            """, unsafe_allow_html=True)
        
        # 示例图片展示
        st.markdown("---")
        st.markdown("### 📸 裂缝识别示例")
        
        col_ex1, col_ex2, col_ex3 = st.columns(3)
        with col_ex1:
            st.markdown("**原始图像**")
            sample_img = create_sample_crack_image()
            st.image(sample_img, caption="混凝土裂缝图像", use_container_width=True)
        with col_ex2:
            st.markdown("**边缘检测**")
            gray = cv2.cvtColor(sample_img, cv2.COLOR_RGB2GRAY)
            edges = cv2.Canny(gray, 30, 90)
            st.image(edges, caption="Canny边缘检测结果", use_container_width=True, clamp=True)
        with col_ex3:
            st.markdown("**检测结果**")
            result = detect_cracks(sample_img, sensitivity=6)
            st.image(result['result_image'], caption="裂缝检测结果（红框标注）", use_container_width=True)
        
        st.caption("💡 说明：裂缝在边缘检测图中表现为连续的白线，CNN通过学习这些特征识别裂缝")
    
    # ==================== 页面2：卷积原理演示 ====================
    elif selected == "🎬 卷积原理演示":
        st.header("🎬 卷积操作原理演示")
        
        tab1, tab2, tab3 = st.tabs(["🔬 卷积过程", "📊 池化过程", "⚡ 激活函数"])
        
        with tab1:
            col1, col2 = st.columns([1.5, 1])
            
            with col1:
                st.markdown("### 🔄 卷积过程可视化")
                st.pyplot(draw_conv_demo())
                
                st.markdown("---")
                st.markdown("### 🎮 互动体验")
                
                kernel_strength = st.slider("卷积核强度调节", 0.0, 2.0, 1.0, 0.1)
                
                fig, ax = plt.subplots(figsize=(10, 3))
                x = np.linspace(0, 10, 100)
                signal = np.sin(x) * 0.5
                signal[45:55] += 1.0 * kernel_strength
                
                ax.plot(x, signal, 'b-', linewidth=2, color='#667eea')
                ax.fill_between(x, 0, signal, alpha=0.3, color='#667eea')
                ax.set_title(f'卷积核强度 = {kernel_strength} 时的特征响应')
                ax.set_xlabel('位置')
                ax.set_ylabel('响应强度')
                ax.grid(True, alpha=0.3)
                st.pyplot(fig)
                st.caption("💡 强度越大，模型对裂缝边缘越敏感")
            
            with col2:
                st.markdown("### 📖 卷积计算原理")
                st.markdown("""
                <div class="card">
                    <h4>卷积三步走：</h4>
                    <ol>
                        <li><strong>滑动</strong> 📍 卷积核在图像上滑动</li>
                        <li><strong>计算</strong> 🧮 对应元素相乘后求和</li>
                        <li><strong>输出</strong> 📤 生成特征图</li>
                    </ol>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown("### 🎯 裂缝识别意义")
                st.info("""
                裂缝在图像中表现为**边缘特征明显**的区域。
                
                - 卷积核可以增强裂缝边缘
                - 抑制背景噪声
                - 突出裂缝特征
                """)
                
                st.markdown("### 🔬 常用卷积核")
                kernel_data = {
                    "边缘检测": "[[-1,-1,-1],\n [-1,8,-1],\n [-1,-1,-1]]",
                    "水平边缘": "[[-1,-2,-1],\n [0,0,0],\n [1,2,1]]",
                    "垂直边缘": "[[-1,0,1],\n [-2,0,2],\n [-1,0,1]]"
                }
                selected_kernel = st.selectbox("选择卷积核类型", list(kernel_data.keys()))
                st.code(kernel_data[selected_kernel], language="python")
        
        with tab2:
            st.markdown("### 📊 池化操作演示")
            st.pyplot(draw_pooling_demo())
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("""
                <div class="card">
                    <h4>最大池化 (Max Pooling)</h4>
                    <ul>
                        <li>取窗口内的最大值</li>
                        <li>保留最显著的特征</li>
                        <li>常用池化方式</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)
            with col2:
                st.markdown("""
                <div class="card">
                    <h4>平均池化 (Average Pooling)</h4>
                    <ul>
                        <li>取窗口内的平均值</li>
                        <li>保留整体信息</li>
                        <li>常用于全连接层前</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)
            
            st.info("💡 **池化的作用**：降维压缩、减少计算量、防止过拟合、增强平移不变性")
        
        with tab3:
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("### ⚡ ReLU激活函数")
                st.pyplot(draw_relu_graph())
            with col2:
                st.markdown("""
                <div class="card">
                    <h4>ReLU 函数特点</h4>
                    <ul>
                        <li>公式：<strong>f(x) = max(0, x)</strong></li>
                        <li>计算简单，速度快</li>
                        <li>缓解梯度消失问题</li>
                        <li>稀疏激活，提高效率</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)
            
            st.markdown("### 📈 常见激活函数对比")
            fig, ax = plt.subplots(figsize=(10, 4))
            x = np.linspace(-3, 3, 100)
            ax.plot(x, np.maximum(0, x), 'r-', label='ReLU', linewidth=2)
            ax.plot(x, 1/(1+np.exp(-x)), 'b-', label='Sigmoid', linewidth=2)
            ax.plot(x, np.tanh(x), 'g-', label='Tanh', linewidth=2)
            ax.axhline(y=0, color='k', linestyle='-', alpha=0.3)
            ax.axvline(x=0, color='k', linestyle='-', alpha=0.3)
            ax.set_xlabel('x')
            ax.set_ylabel('f(x)')
            ax.set_title('激活函数对比')
            ax.legend()
            ax.grid(True, alpha=0.3)
            st.pyplot(fig)
    
    # ==================== 页面3：裂缝检测实战 ====================
    elif selected == "🔍 裂缝检测实战":
        st.header("🔍 裂缝检测实战演示")
        st.markdown("上传混凝土图片，系统自动识别**黑色细长裂缝**并用红色矩形框标注")
        
        col_left, col_right = st.columns([1.2, 0.8])
        
        with col_left:
            st.subheader("📸 输入图像")
            
            uploaded_file = st.file_uploader("上传混凝土图片", type=['jpg', 'png', 'jpeg'])
            
            # 示例图片快捷选择
            st.markdown("**📋 快速测试：**")
            col_sample1, col_sample2 = st.columns(2)
            with col_sample1:
                if st.button("🔴 测试：有裂缝", use_container_width=True):
                    st.session_state.test_image = create_sample_crack_image()
                    st.session_state.has_test_image = True
            with col_sample2:
                if st.button("✅ 测试：正常混凝土", use_container_width=True):
                    st.session_state.test_image = create_sample_normal_image()
                    st.session_state.has_test_image = True
            
            if uploaded_file:
                image = Image.open(uploaded_file)
                st.image(image, caption="待检测图片", use_container_width=True)
                st.caption(f"📷 图片尺寸: {image.size[0]} × {image.size[1]} 像素")
            elif st.session_state.get('has_test_image', False):
                image = Image.fromarray(st.session_state.test_image)
                st.image(image, caption="测试图片", use_container_width=True)
            else:
                st.info("👈 请上传图片或点击测试按钮")
                
                # 展示示例效果
                st.markdown("---")
                st.markdown("### 📷 检测效果示例")
                sample = create_sample_crack_image()
                result = detect_cracks(sample)
                st.image(result['result_image'], caption="裂缝检测示例（红色框为结果）", use_container_width=True)
        
        with col_right:
            st.subheader("⚙️ 检测参数")
            
            sensitivity = st.slider("🔍 灵敏度", 1, 10, 6, help="1=只检测明显裂缝，10=检测所有可能裂缝")
            min_length = st.slider("📏 最小裂缝长度", 20, 200, 50, help="长度小于此值的区域将被忽略")
            connect_distance = st.slider("🔗 连接距离", 5, 50, 20, help="断裂裂缝的最大连接距离")
            
            st.markdown("---")
            
            if sensitivity <= 3:
                st.info("🔴 **低灵敏度模式**：只检测非常明显的裂缝")
            elif sensitivity <= 7:
                st.success("🟡 **中灵敏度模式**：检测常规裂缝")
            else:
                st.warning("🟢 **高灵敏度模式**：检测细小裂缝（可能有误报）")
            
            # 检测按钮
            image_to_process = None
            if uploaded_file:
                image_to_process = Image.open(uploaded_file)
            elif st.session_state.get('has_test_image', False):
                image_to_process = Image.fromarray(st.session_state.test_image)
            
            if image_to_process:
                if st.button("🔍 开始检测", type="primary", use_container_width=True):
                    with st.spinner("🔬 正在分析图像..."):
                        result = detect_cracks(image_to_process, sensitivity, min_length, connect_distance)
                        
                        # 保存记录
                        conn = sqlite3.connect('demo_records.db')
                        conn.execute(
                            'INSERT INTO records (session_id, image_name, crack_result, confidence, crack_count, timestamp, parameters) VALUES (?,?,?,?,?,?,?)',
                            (get_session_id(), uploaded_file.name if uploaded_file else "测试图片",
                             "裂缝" if result['has_crack'] else "正常",
                             result['confidence'], result['crack_count'],
                             datetime.now().isoformat(),
                             f"sens={sensitivity},len={min_length},conn={connect_distance}")
                        )
                        conn.commit()
                        conn.close()
                        
                        st.markdown("### 📊 检测结果")
                        
                        if result['has_crack']:
                            st.error(f"⚠️ 检测到 {result['crack_count']} 处裂缝")
                            st.metric("置信度", f"{result['confidence']:.1%}")
                            if result['crack_info']:
                                st.caption(f"裂缝总面积: {sum(c['area'] for c in result['crack_info']):.0f} 像素")
                        else:
                            st.success("✅ 未检测到裂缝")
                            if sensitivity < 8:
                                st.info("💡 提示：可以尝试调高「灵敏度」或调低「最小长度」")
                        
                        st.image(result['result_image'], caption="检测结果（红色矩形框为裂缝位置）", use_container_width=True)
                        st.caption(f"🆔 会话ID: {get_session_id()}")
            else:
                st.info("👈 请先上传图片或点击测试按钮")
    
    # ==================== 页面4：智能测验 ====================
    elif selected == "📝 智能测验":
        st.header("📝 CNN知识智能测验")
        st.markdown("完成以下10道测验题，系统将智能分析你的学习情况并提供个性化建议")
        
        if 'quiz_answers' not in st.session_state:
            st.session_state.quiz_answers = {}
        
        # 显示进度
        completed = len(st.session_state.quiz_answers)
        st.progress(completed / len(QUESTION_BANK))
        st.caption(f"已完成 {completed}/{len(QUESTION_BANK)} 题")
        
        for i, q in enumerate(QUESTION_BANK):
            with st.container():
                st.markdown(f"""
                <div class="quiz-card">
                    <h4>📌 第{i+1}题 <span style="font-size:0.8rem; color:#667eea">({q['category']})</span></h4>
                    <p style="font-size:1.1rem">{q['question']}</p>
                </div>
                """, unsafe_allow_html=True)
                
                user_answer = st.radio(
                    "请选择答案",
                    q['options'],
                    key=f"quiz_{q['id']}",
                    label_visibility="collapsed",
                    index=None
                )
                
                if user_answer:
                    st.session_state.quiz_answers[q['id']] = {
                        'question': q['question'],
                        'user_answer': user_answer,
                        'correct_answer': q['correct_text'],
                        'is_correct': user_answer.startswith(q['correct']),
                        'explanation': q['explanation'],
                        'category': q['category']
                    }
            
            st.markdown("---")
        
        col1, col2, col3 = st.columns([1, 1, 1])
        with col2:
            if st.button("📊 提交并获取智能分析", type="primary", use_container_width=True):
                if len(st.session_state.quiz_answers) < len(QUESTION_BANK):
                    st.warning(f"请先完成所有题目（已完成 {len(st.session_state.quiz_answers)}/{len(QUESTION_BANK)}）")
                else:
                    correct_count = sum(1 for a in st.session_state.quiz_answers.values() if a['is_correct'])
                    total_score = int(correct_count / len(QUESTION_BANK) * 100)
                    
                    conn = sqlite3.connect('demo_records.db')
                    for qid, answer in st.session_state.quiz_answers.items():
                        conn.execute(
                            'INSERT INTO quiz_records (session_id, question_id, question, user_answer, correct_answer, is_correct, category, timestamp) VALUES (?,?,?,?,?,?,?,?)',
                            (get_session_id(), qid, answer['question'], answer['user_answer'], 
                             answer['correct_answer'], 1 if answer['is_correct'] else 0,
                             answer['category'], datetime.now().isoformat())
                        )
                    conn.commit()
                    conn.close()
                    
                    st.session_state.show_analysis = True
                    st.session_state.quiz_result = {
                        'correct_count': correct_count,
                        'total': len(QUESTION_BANK),
                        'score': total_score,
                        'answers': list(st.session_state.quiz_answers.values())
                    }
                    st.rerun()
        
        if st.session_state.get('show_analysis', False):
            st.markdown("---")
            st.markdown("## 🤖 智能学情分析报告")
            
            result = st.session_state.quiz_result
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("正确题数", f"{result['correct_count']}/{result['total']}")
            with col2:
                st.metric("得分", f"{result['score']}分")
            with col3:
                if result['score'] >= 80:
                    st.metric("评级", "卓越 🏆")
                elif result['score'] >= 60:
                    st.metric("评级", "良好 📈")
                else:
                    st.metric("评级", "需加强 💪")
            
            # 雷达图
            topic_stats = {}
            for a in result['answers']:
                topic = a['category']
                if topic not in topic_stats:
                    topic_stats[topic] = {'correct': 0, 'total': 0}
                topic_stats[topic]['total'] += 1
                if a['is_correct']:
                    topic_stats[topic]['correct'] += 1
            
            if topic_stats:
                fig = go.Figure(data=go.Scatterpolar(
                    r=[stats['correct']/stats['total']*100 for stats in topic_stats.values()],
                    theta=list(topic_stats.keys()),
                    fill='toself',
                    marker=dict(color='#667eea')
                ))
                fig.update_layout(polar=dict(radialaxis=dict(visible=True, range=[0, 100])), title="各知识点掌握度")
                st.plotly_chart(fig, use_container_width=True)
            
            with st.spinner("正在生成智能分析报告..."):
                analysis = generate_smart_analysis(
                    result['answers'],
                    get_session_id(),
                    result['score']
                )
            
            st.markdown(analysis, unsafe_allow_html=True)
            
            with st.expander("📋 查看错题详解"):
                wrong_count = 0
                for i, a in enumerate(result['answers']):
                    if not a['is_correct']:
                        wrong_count += 1
                        st.markdown(f"""
                        **第{i+1}题** - {a['category']}
                        - 题目：{a['question']}
                        - 你的答案：{a['user_answer']}
                        - 正确答案：{a['correct_answer']}
                        - 📖 解析：{a['explanation']}
                        ---
                        """)
                if wrong_count == 0:
                    st.success("🎉 太棒了！全对！继续保持！")
            
            if st.button("🔄 重新测验", use_container_width=True):
                st.session_state.quiz_answers = {}
                st.session_state.show_analysis = False
                st.rerun()
    
    # ==================== 页面5：学情分析 ====================
    else:
        st.header("📊 学情分析看板")
        
        stats = get_quiz_statistics(get_session_id())
        
        if stats and stats['total'] > 0:
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("已完成测验", f"{stats['total']}题")
            with col2:
                st.metric("正确题数", f"{stats['correct']}")
            with col3:
                st.metric("平均得分", f"{stats['score']}分")
            
            if stats['topic_stats']:
                df_topic = pd.DataFrame([
                    {'知识点': t, '掌握度': v['correct']/v['total']*100}
                    for t, v in stats['topic_stats'].items()
                ])
                fig = px.bar(df_topic, x='知识点', y='掌握度', title='各知识点掌握度',
                            color='掌握度', color_continuous_scale=['#e74c3c', '#f39c12', '#2ecc71'],
                            range_y=[0, 100])
                st.plotly_chart(fig, use_container_width=True)
            
            if len(stats['records']) > 1:
                df_time = pd.DataFrame(stats['records'])
                df_time['timestamp'] = pd.to_datetime(df_time['timestamp'])
                df_time = df_time.sort_values('timestamp')
                df_time['累积正确率'] = df_time['is_correct'].expanding().mean() * 100
                
                fig = px.line(df_time, x='timestamp', y='累积正确率', title='学习进步曲线')
                st.plotly_chart(fig, use_container_width=True)
            
            conn = sqlite3.connect('demo_records.db')
            analysis_record = conn.execute(
                'SELECT analysis_text, score, timestamp FROM analysis_records WHERE session_id = ? ORDER BY id DESC LIMIT 1',
                (get_session_id(),)
            ).fetchone()
            conn.close()
            
            if analysis_record:
                with st.expander("📋 最新分析报告", expanded=True):
                    st.caption(f"分析时间：{analysis_record[2]} | 得分：{analysis_record[1]}分")
                    st.markdown(analysis_record[0], unsafe_allow_html=True)
        else:
            st.info("📊 暂无学情数据，请先完成「智能测验」")
            if st.button("🎯 前往智能测验"):
                st.session_state.nav_menu = "📝 智能测验"
                st.rerun()
    
    st.markdown("---")
    st.caption("💡 提示：完成智能测验后，系统会生成个性化学习建议 | 裂缝检测用红色矩形框标注")

if __name__ == "__main__":
    main()